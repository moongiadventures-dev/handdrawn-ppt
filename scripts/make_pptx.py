#!/usr/bin/env python3
"""
handdrawn-ppt — 손그림 삽화 + 텍스트를 실제 .pptx 파일로 조립한다.

캐러셀(4:5 PNG)과는 별개다. 이쪽은 발표장에서 여는 16:9 슬라이드.

사용법
------
    python make_pptx.py deck.json -o 발표자료.pptx

deck.json 스펙
--------------
{
  "title": "AI와 내 일자리",              // 파일 메타용 (선택)
  "footer": "MOONGI Studio",              // 모든 슬라이드 우하단 (선택)
  "font": "Gaegu",                        // 손글씨 폰트 이름 (선택, 기본 Gaegu)
  "body_font": "Noto Sans KR",            // 출처·노트용 (선택)
  "slides": [
    {
      "layout": "title",                  // title | content | full | quote
      "tag": "손그림 PPT",                 // 좌상단 작은 라벨 (선택)
      "headline": "내 캐릭터 손그림으로\\nPPT 자료 자동화하기",
      "body": "설명 줄1\\n설명 줄2",        // 선택
      "image": "illustrations/01.png",    // deck.json 기준 상대경로 (선택)
      "source": "출처 한 줄",              // 선택 — 슬라이드 하단 + 발표자 노트
      "notes": "발표할 때 할 말"           // 선택 — 발표자 노트에 추가
    }
  ]
}

레이아웃
--------
  title    — 헤드라인 크게 왼쪽, 삽화 오른쪽
  content  — 헤드라인 위, 본문 왼쪽, 삽화 오른쪽   (기본값)
  full     — 헤드라인 위, 삽화 가운데 크게, 본문 아래
  quote    — 텍스트만, 가운데 정렬 (전환 슬라이드용)

폰트 주의
---------
pptx는 폰트 '이름'만 저장한다. 파일을 여는 컴퓨터에 그 폰트가 없으면 다른 걸로 대체된다.
손글씨 느낌을 유지하려면 발표할 PC에 Gaegu를 설치해야 한다.
  https://fonts.google.com/specimen/Gaegu
설치가 어려우면 deck.json에서 "font"를 "맑은 고딕" 같은 기본 폰트로 바꿔도 동작은 한다.
"""

import argparse
import json
import sys
from pathlib import Path

# 윈도우 기본 콘솔(cp949/cp1252)은 한글 경로를 print하다 죽는다.
# 파일 저장은 이미 끝난 뒤라 더 억울하므로 출력 인코딩을 먼저 고정한다.
for _stream in (sys.stdout, sys.stderr):
    try:
        _stream.reconfigure(encoding="utf-8", errors="replace")
    except Exception:
        pass

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Emu, Inches, Pt
except ImportError:
    sys.exit("python-pptx가 없다.  pip install python-pptx")

try:
    from PIL import Image
except ImportError:
    Image = None  # 이미지 비율 계산만 못 함. 없어도 동작은 한다.

# ── 슬라이드 기본값 (16:9) ────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.62)

INK = RGBColor(0x1A, 0x1A, 0x1A)
ACCENT = RGBColor(0xE2, 0x4A, 0x3B)
GREY = RGBColor(0x9A, 0x9A, 0x9A)
BODY_GREY = RGBColor(0x33, 0x33, 0x33)

DEFAULT_FONT = "Gaegu"
DEFAULT_BODY_FONT = "Noto Sans KR"


def _tb(slide, x, y, w, h):
    """여백 0짜리 텍스트박스."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _write(tf, text, *, font, size, color, bold=False, line=1.25,
           align=PP_ALIGN.LEFT, space_after=0):
    """여러 줄 텍스트를 문단별로 채운다. 첫 문단은 재사용."""
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line
        para.space_after = Pt(space_after)
        run = para.add_run()
        run.text = ln
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def _fit(img_path, box_w, box_h):
    """박스 안에 비율 유지로 들어가는 (w, h). PIL 없으면 16:9 가정."""
    ratio = 16 / 9
    if Image is not None:
        try:
            with Image.open(img_path) as im:
                ratio = im.width / im.height
        except Exception:
            pass
    if box_w / box_h > ratio:          # 박스가 더 납작 → 높이 기준
        h = box_h
        w = Emu(int(box_h * ratio))
    else:                              # 박스가 더 길쭉 → 폭 기준
        w = box_w
        h = Emu(int(box_w / ratio))
    return w, h


def _shrink(img_path, max_w, cache_dir):
    """
    2K~4K 원본을 그대로 넣으면 pptx가 수십 MB가 된다.
    슬라이드에 실제로 보이는 크기는 1600px 남짓이라 그 이상은 낭비다.
    max_w보다 크면 축소본을 만들어 그 경로를 돌려준다.
    """
    if Image is None or not max_w:
        return img_path
    try:
        with Image.open(img_path) as im:
            if im.width <= max_w:
                return img_path
            ratio = max_w / im.width
            out = cache_dir / f"{img_path.stem}__{max_w}.png"
            if not out.exists():
                cache_dir.mkdir(parents=True, exist_ok=True)
                im.convert("RGB").resize(
                    (max_w, int(im.height * ratio)), Image.LANCZOS
                ).save(out, "PNG", optimize=True)
            return out
    except Exception:
        return img_path


def _place_image(slide, img_path, x, y, box_w, box_h):
    """박스 중앙에 비율 유지로 배치."""
    w, h = _fit(img_path, box_w, box_h)
    slide.shapes.add_picture(
        str(img_path),
        Emu(int(x + (box_w - w) / 2)),
        Emu(int(y + (box_h - h) / 2)),
        width=w, height=h,
    )


def _add_tag(slide, text, font):
    if not text:
        return
    tf = _tb(slide, MARGIN, Inches(0.42), Inches(6), Inches(0.4))
    _write(tf, text, font=font, size=15, color=ACCENT, bold=True)


def _add_source(slide, text, body_font):
    if not text:
        return
    tf = _tb(slide, MARGIN, Inches(6.72), SLIDE_W - MARGIN * 2, Inches(0.5))
    _write(tf, text, font=body_font, size=9, color=GREY, line=1.35)


def _add_footer(slide, text, body_font):
    if not text:
        return
    tf = _tb(slide, SLIDE_W - Inches(4.2), Inches(6.86),
             Inches(3.6), Inches(0.32))
    _write(tf, text, font=body_font, size=9, color=GREY, align=PP_ALIGN.RIGHT)


def _add_notes(slide, source, notes):
    parts = [p for p in (notes, (f"출처: {source}" if source else None)) if p]
    if parts:
        slide.notes_slide.notes_text_frame.text = "\n\n".join(parts)


# ── 레이아웃 ──────────────────────────────────────────────────────

def layout_title(slide, s, font, body_font):
    tf = _tb(slide, MARGIN, Inches(2.05), Inches(6.1), Inches(3.0))
    _write(tf, s["headline"], font=font, size=48, color=INK, bold=True, line=1.12)
    if s.get("body"):
        tf2 = _tb(slide, MARGIN, Inches(5.05), Inches(6.1), Inches(1.2))
        _write(tf2, s["body"], font=font, size=19, color=BODY_GREY, line=1.4)
    if s.get("_img"):
        _place_image(slide, s["_img"], Inches(6.95), Inches(1.35),
                     Inches(5.75), Inches(4.6))


def layout_content(slide, s, font, body_font):
    tf = _tb(slide, MARGIN, Inches(1.0), Inches(11.9), Inches(1.3))
    _write(tf, s["headline"], font=font, size=36, color=INK, bold=True, line=1.1)

    head_lines = len(s["headline"].split("\n"))
    top = Inches(1.0) + Inches(0.55) * head_lines + Inches(0.3)

    if s.get("_img"):
        _place_image(slide, s["_img"], Inches(5.55), top,
                     Inches(7.15), Inches(6.55) - top)
        body_w = Inches(4.6)
    else:
        body_w = Inches(11.9)

    if s.get("body"):
        tf2 = _tb(slide, MARGIN, top + Inches(0.12), body_w, Inches(3.6))
        _write(tf2, s["body"], font=font, size=19, color=BODY_GREY,
               line=1.45, space_after=7)


def layout_full(slide, s, font, body_font):
    tf = _tb(slide, MARGIN, Inches(0.95), Inches(11.9), Inches(1.0))
    _write(tf, s["headline"], font=font, size=36, color=INK, bold=True, line=1.1)

    body_h = Inches(0.95) if s.get("body") else Inches(0)
    if s.get("_img"):
        _place_image(slide, s["_img"], MARGIN, Inches(2.0),
                     SLIDE_W - MARGIN * 2, Inches(4.5) - body_h)
    if s.get("body"):
        tf2 = _tb(slide, MARGIN, Inches(5.72), SLIDE_W - MARGIN * 2, Inches(0.95))
        _write(tf2, s["body"], font=font, size=18, color=BODY_GREY,
               line=1.4, align=PP_ALIGN.CENTER)


def layout_quote(slide, s, font, body_font):
    tf = _tb(slide, Inches(1.6), Inches(2.5), Inches(10.1), Inches(2.4))
    _write(tf, s["headline"], font=font, size=44, color=INK, bold=True,
           line=1.18, align=PP_ALIGN.CENTER)
    if s.get("body"):
        tf2 = _tb(slide, Inches(1.6), Inches(4.85), Inches(10.1), Inches(1.2))
        _write(tf2, s["body"], font=font, size=19, color=BODY_GREY,
               line=1.4, align=PP_ALIGN.CENTER)


LAYOUTS = {
    "title": layout_title,
    "content": layout_content,
    "full": layout_full,
    "quote": layout_quote,
}


# ── 빌드 ──────────────────────────────────────────────────────────

def build(spec_path: Path, out_path: Path, max_img_width: int = 1800) -> Path:
    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    base = spec_path.parent

    font = spec.get("font", DEFAULT_FONT)
    body_font = spec.get("body_font", DEFAULT_BODY_FONT)
    footer = spec.get("footer", "")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    warnings = []

    for i, s in enumerate(spec["slides"], 1):
        if not s.get("headline"):
            warnings.append(f"slide {i}: headline 없음")
            s["headline"] = ""

        s["_img"] = None
        if s.get("image"):
            p = (base / s["image"]).resolve()
            if p.exists():
                s["_img"] = _shrink(p, max_img_width, base / ".pptx_cache")
            else:
                warnings.append(f"slide {i}: 이미지 없음 -> {s['image']}")

        slide = prs.slides.add_slide(blank)
        kind = s.get("layout", "content")
        if kind not in LAYOUTS:
            warnings.append(f"slide {i}: 모르는 layout '{kind}' -> content로 처리")
            kind = "content"

        _add_tag(slide, s.get("tag"), font)
        LAYOUTS[kind](slide, s, font, body_font)
        _add_source(slide, s.get("source"), body_font)
        _add_footer(slide, footer, body_font)
        _add_notes(slide, s.get("source"), s.get("notes"))

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    size_mb = out_path.stat().st_size / (1024 * 1024)
    print(f"built: {out_path}  ({len(spec['slides'])} slides, 16:9, {size_mb:.1f} MB)")
    if warnings:
        print("\n[warnings]")
        for w in warnings:
            print("  -", w)
    print(f"\nfont: '{font}' — 여는 PC에 없으면 대체 폰트로 표시된다.")
    print("      https://fonts.google.com/specimen/Gaegu")
    return out_path


def main():
    ap = argparse.ArgumentParser(
        description="손그림 삽화 + 텍스트 -> 16:9 .pptx")
    ap.add_argument("spec", help="deck.json 경로")
    ap.add_argument("-o", "--output", help="출력 .pptx 경로")
    ap.add_argument("--max-image-width", type=int, default=1800,
                    metavar="PX",
                    help="이미지를 이 폭으로 축소해서 넣는다 (기본 1800). "
                         "0이면 원본 그대로 — 파일이 수십 MB가 될 수 있다.")
    args = ap.parse_args()

    spec_path = Path(args.spec).resolve()
    if not spec_path.exists():
        sys.exit(f"스펙 파일이 없다: {spec_path}")

    out = Path(args.output) if args.output else spec_path.with_suffix(".pptx")
    build(spec_path, out.resolve(), max_img_width=args.max_image_width)


if __name__ == "__main__":
    main()
